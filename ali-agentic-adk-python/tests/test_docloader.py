from __future__ import annotations

import io
import os
import json
from pathlib import Path
from typing import Any, Dict, Iterable, Tuple

import pytest
from PyPDF2 import PdfWriter
from docx import Document as WordDocument
import openpyxl
from pptx import Presentation
from pptx.util import Inches

from ali_agentic_adk_python.core import (
    Document,
    ExcelDocLoader,
    FeishuDocLoader,
    MarkdownDocLoader,
    PDFDocLoader,
    PPTDocLoader,
    RecursiveCharacterTextSplitter,
    TextDocLoader,
    WordDocLoader,
    YuqueDocLoader,
)


class _StubHTTPResponse:
    def __init__(self, payload: Any):
        if isinstance(payload, tuple):
            body, status = payload
        else:
            body, status = payload, 200
        self._body = body
        self.status_code = status

    def json(self) -> Any:
        return self._body

    def raise_for_status(self) -> None:
        if self.status_code >= 400:
            raise RuntimeError(f"HTTP {self.status_code}")


class _StubSession:
    def __init__(self, responses: Dict[Tuple[str, str], Iterable[Any]]):
        self._responses: Dict[Tuple[str, str], Iterable[Any]] = {}
        for key, value in responses.items():
            if isinstance(value, list):
                self._responses[key] = iter(value)
            else:
                self._responses[key] = iter([value])
        self.calls: list[Tuple[str, str, Dict[str, Any]]] = []

    def request(self, method: str, url: str, **kwargs: Any) -> _StubHTTPResponse:
        key = (method.upper(), url)
        if key not in self._responses:
            raise AssertionError(f"Unexpected request: {key}")
        self.calls.append((method.upper(), url, kwargs))
        iterator = self._responses[key]
        try:
            payload = next(iterator)
        except StopIteration:
            raise AssertionError(f"No more responses configured for {key}")
        return _StubHTTPResponse(payload)

    def get(self, url: str, **kwargs: Any) -> _StubHTTPResponse:
        return self.request("GET", url, **kwargs)

    def post(self, url: str, **kwargs: Any) -> _StubHTTPResponse:
        return self.request("POST", url, **kwargs)

# ============================================================================
# TextDocLoader 基础测试
# ============================================================================

def test_text_doc_loader_reads_file(tmp_path):
    path = tmp_path / "sample.txt"
    path.write_text("Hello ADK!", encoding="utf-8")

    loader = TextDocLoader(file_path=str(path))
    documents = loader.load()

    assert len(documents) == 1
    doc = documents[0]
    assert doc.page_content == "Hello ADK!"
    assert doc.metadata["source"] == str(path)


def test_text_doc_loader_raises_when_path_missing():
    loader = TextDocLoader()
    with pytest.raises(ValueError):
        loader.load()


def test_text_doc_loader_fetches_from_metadata_stream():
    loader = TextDocLoader()
    stream = io.StringIO("From stream")

    documents = loader.fetch_content({"stream": stream, "metadata": {"category": "demo"}})

    assert documents[0].page_content == "From stream"
    assert documents[0].metadata["source"] == "stream"
    assert documents[0].metadata["category"] == "demo"


def test_text_doc_loader_fetches_from_metadata_path(tmp_path):
    path = tmp_path / "note.txt"
    path.write_text("metadata path", encoding="utf-8")

    loader = TextDocLoader()
    documents = loader.load_with_meta({"file_path": str(path)})

    assert documents[0].page_content == "metadata path"
    assert documents[0].metadata["source"] == str(path)


def test_text_doc_loader_load_and_split_preserves_metadata(tmp_path):
    path = tmp_path / "long.txt"
    content = "alpha beta gamma delta epsilon zeta eta theta"
    path.write_text(content, encoding="utf-8")

    loader = TextDocLoader(str(path))
    splitter = RecursiveCharacterTextSplitter(max_chunk_size=15, max_chunk_overlap=3)

    chunks = loader.load_and_split(splitter)

    assert len(chunks) > 1
    assert {chunk.metadata["source"] for chunk in chunks} == {str(path)}
    combined = "".join(chunk.page_content for chunk in chunks).replace(" ", "")
    assert combined == content.replace(" ", "")


# ============================================================================
# TextDocLoader 编码测试
# ============================================================================

def test_text_doc_loader_reads_utf8_with_bom(tmp_path):
    """测试读取带BOM的UTF-8文件"""
    path = tmp_path / "utf8_bom.txt"
    content = "UTF-8 with BOM: 你好世界"
    path.write_bytes(b'\xef\xbb\xbf' + content.encode('utf-8'))
    
    loader = TextDocLoader(file_path=str(path))
    documents = loader.load()
    
    assert len(documents) == 1
    assert "你好世界" in documents[0].page_content


def test_text_doc_loader_reads_chinese_characters(tmp_path):
    """测试读取中文字符"""
    path = tmp_path / "chinese.txt"
    content = "这是一个中文测试文件，包含中文字符。"
    path.write_text(content, encoding="utf-8")
    
    loader = TextDocLoader(file_path=str(path))
    documents = loader.load()
    
    assert documents[0].page_content == content
    assert "中文" in documents[0].page_content


def test_text_doc_loader_reads_emoji(tmp_path):
    """测试读取emoji表情"""
    path = tmp_path / "emoji.txt"
    content = "Hello 👋 World 🌍 ADK 🚀"
    path.write_text(content, encoding="utf-8")
    
    loader = TextDocLoader(file_path=str(path))
    documents = loader.load()
    
    assert documents[0].page_content == content


def test_text_doc_loader_reads_mixed_languages(tmp_path):
    """测试读取混合语言文本"""
    path = tmp_path / "mixed.txt"
    content = "English 英文 日本語 한국어 Español Français"
    path.write_text(content, encoding="utf-8")
    
    loader = TextDocLoader(file_path=str(path))
    documents = loader.load()
    
    assert documents[0].page_content == content


# ============================================================================
# TextDocLoader 空内容和特殊内容测试
# ============================================================================

def test_text_doc_loader_reads_empty_file(tmp_path):
    """测试读取空文件"""
    path = tmp_path / "empty.txt"
    path.write_text("", encoding="utf-8")
    
    loader = TextDocLoader(file_path=str(path))
    documents = loader.load()
    
    assert len(documents) == 1
    assert documents[0].page_content == ""


def test_text_doc_loader_reads_whitespace_only(tmp_path):
    """测试只包含空白字符的文件"""
    path = tmp_path / "whitespace.txt"
    path.write_text("   \n\n   \t\t   \n", encoding="utf-8")
    
    loader = TextDocLoader(file_path=str(path))
    documents = loader.load()
    
    assert len(documents) == 1
    assert documents[0].page_content.strip() == ""


def test_text_doc_loader_reads_single_line(tmp_path):
    """测试单行文本"""
    path = tmp_path / "single_line.txt"
    path.write_text("Single line without newline", encoding="utf-8")
    
    loader = TextDocLoader(file_path=str(path))
    documents = loader.load()
    
    assert documents[0].page_content == "Single line without newline"


def test_text_doc_loader_reads_multiline(tmp_path):
    """测试多行文本"""
    path = tmp_path / "multiline.txt"
    content = "Line 1\nLine 2\nLine 3\nLine 4\nLine 5"
    path.write_text(content, encoding="utf-8")
    
    loader = TextDocLoader(file_path=str(path))
    documents = loader.load()
    
    assert documents[0].page_content == content
    assert documents[0].page_content.count("\n") == 4


def test_text_doc_loader_reads_long_lines(tmp_path):
    """测试超长行"""
    path = tmp_path / "long_line.txt"
    content = "x" * 10000
    path.write_text(content, encoding="utf-8")
    
    loader = TextDocLoader(file_path=str(path))
    documents = loader.load()
    
    assert len(documents[0].page_content) == 10000


def test_text_doc_loader_reads_special_characters(tmp_path):
    """测试特殊字符"""
    path = tmp_path / "special.txt"
    content = "Special chars: !@#$%^&*()_+-=[]{}|;:',.<>?/~`"
    path.write_text(content, encoding="utf-8")
    
    loader = TextDocLoader(file_path=str(path))
    documents = loader.load()
    
    assert documents[0].page_content == content


def test_text_doc_loader_reads_json_content(tmp_path):
    """测试JSON格式内容"""
    path = tmp_path / "data.json"
    data = {"name": "test", "value": 123, "nested": {"key": "value"}}
    path.write_text(json.dumps(data, indent=2), encoding="utf-8")
    
    loader = TextDocLoader(file_path=str(path))
    documents = loader.load()
    
    loaded_data = json.loads(documents[0].page_content)
    assert loaded_data == data


# ============================================================================
# TextDocLoader 路径和文件系统测试
# ============================================================================

def test_text_doc_loader_with_relative_path(tmp_path):
    """测试相对路径"""
    old_cwd = os.getcwd()
    try:
        os.chdir(tmp_path)
        path = Path("relative.txt")
        path.write_text("Relative path test", encoding="utf-8")
        
        loader = TextDocLoader(file_path="relative.txt")
        documents = loader.load()
        
        assert documents[0].page_content == "Relative path test"
    finally:
        os.chdir(old_cwd)


def test_text_doc_loader_with_nested_directories(tmp_path):
    """测试嵌套目录"""
    nested_dir = tmp_path / "level1" / "level2" / "level3"
    nested_dir.mkdir(parents=True)
    path = nested_dir / "nested.txt"
    path.write_text("Nested file", encoding="utf-8")
    
    loader = TextDocLoader(file_path=str(path))
    documents = loader.load()
    
    assert documents[0].page_content == "Nested file"


def test_text_doc_loader_nonexistent_file():
    """测试不存在的文件"""
    loader = TextDocLoader(file_path="/nonexistent/path/file.txt")
    with pytest.raises((FileNotFoundError, OSError)):
        loader.load()


# ============================================================================
# TextDocLoader Stream 测试
# ============================================================================

def test_text_doc_loader_stream_with_multiple_lines():
    """测试流读取多行内容"""
    loader = TextDocLoader()
    stream = io.StringIO("Line 1\nLine 2\nLine 3")
    
    documents = loader.fetch_content({"stream": stream})
    
    assert documents[0].page_content == "Line 1\nLine 2\nLine 3"


def test_text_doc_loader_stream_empty():
    """测试空流"""
    loader = TextDocLoader()
    stream = io.StringIO("")
    
    documents = loader.fetch_content({"stream": stream})
    
    assert documents[0].page_content == ""


def test_text_doc_loader_stream_with_custom_metadata():
    """测试带自定义元数据的流"""
    loader = TextDocLoader()
    stream = io.StringIO("Test content")
    metadata = {
        "source": "custom_source",
        "author": "test_author",
        "timestamp": "2025-01-01",
        "category": "test"
    }
    
    documents = loader.fetch_content({"stream": stream, "metadata": metadata})
    
    assert documents[0].metadata["author"] == "test_author"
    assert documents[0].metadata["timestamp"] == "2025-01-01"
    assert documents[0].metadata["category"] == "test"


# ============================================================================
# Document 类测试
# ============================================================================

def test_document_helpers():
    doc = Document(page_content="data", metadata={"category": "notes"}, category="notes")
    assert doc.has_metadata() is True
    assert doc.has_category() is True
    cloned = doc.copy_with_page_content("new")
    assert cloned.page_content == "new"
    assert cloned.metadata == doc.metadata
    assert cloned.metadata is not doc.metadata
    assert cloned.category == doc.category


def test_document_creation_basic():
    """测试基本文档创建"""
    doc = Document(page_content="test content")
    assert doc.page_content == "test content"
    assert doc.metadata == {}
    assert doc.category is None


def test_document_creation_with_metadata():
    """测试带元数据的文档创建"""
    metadata = {"source": "test.txt", "page": 1}
    doc = Document(page_content="content", metadata=metadata)
    
    assert doc.metadata["source"] == "test.txt"
    assert doc.metadata["page"] == 1


def test_document_creation_with_category():
    """测试带分类的文档创建"""
    doc = Document(page_content="content", category="technical")
    assert doc.category == "technical"
    assert doc.has_category() is True


def test_document_empty_content():
    """测试空内容文档"""
    doc = Document(page_content="")
    assert doc.page_content == ""
    assert len(doc.page_content) == 0


def test_document_metadata_modification():
    """测试元数据修改"""
    doc = Document(page_content="content", metadata={"key": "value"})
    doc.metadata["new_key"] = "new_value"
    
    assert doc.metadata["key"] == "value"
    assert doc.metadata["new_key"] == "new_value"


def test_document_copy_independence():
    """测试文档复制的独立性"""
    original = Document(
        page_content="original",
        metadata={"key": "value"},
        category="cat1"
    )
    copy = original.copy_with_page_content("modified")
    
    copy.metadata["key"] = "modified_value"
    copy.metadata["new_key"] = "new"
    
    assert original.metadata["key"] == "value"
    assert "new_key" not in original.metadata
    assert copy.page_content == "modified"
    assert original.page_content == "original"


def test_document_has_metadata_false():
    """测试没有元数据的情况"""
    doc = Document(page_content="content")
    result = doc.has_metadata()
    assert result is not True or len(doc.metadata) == 0


def test_document_has_category_false():
    """测试没有分类的情况"""
    doc = Document(page_content="content")
    assert doc.has_category() is False or doc.category is None


def test_document_long_content():
    """测试长文档内容"""
    long_content = "test " * 10000
    doc = Document(page_content=long_content)
    assert len(doc.page_content) == len(long_content)


def test_document_with_complex_metadata():
    """测试复杂元数据"""
    metadata = {
        "source": "test.txt",
        "nested": {"key1": "value1", "key2": "value2"},
        "list": [1, 2, 3, 4, 5],
        "number": 123,
        "float": 3.14,
        "bool": True
    }
    doc = Document(page_content="content", metadata=metadata)
    
    assert doc.metadata["nested"]["key1"] == "value1"
    assert doc.metadata["list"] == [1, 2, 3, 4, 5]
    assert doc.metadata["number"] == 123
    assert doc.metadata["float"] == 3.14
    assert doc.metadata["bool"] is True


# ============================================================================
# RecursiveCharacterTextSplitter 测试
# ============================================================================

def test_splitter_basic_split():
    """测试基本分割功能"""
    splitter = RecursiveCharacterTextSplitter(max_chunk_size=10, max_chunk_overlap=2)
    text = "This is a test text for splitting"
    chunks = splitter.split_text(text)
    
    assert len(chunks) > 1
    for chunk in chunks:
        assert len(chunk) <= 10 or " " not in chunk[:10]


def test_splitter_no_overlap():
    """测试无重叠分割"""
    splitter = RecursiveCharacterTextSplitter(max_chunk_size=20, max_chunk_overlap=0)
    text = "word1 word2 word3 word4 word5"
    chunks = splitter.split_text(text)
    
    combined = "".join(chunks)
    assert combined.replace(" ", "") == text.replace(" ", "")


def test_splitter_with_overlap():
    """测试有重叠分割"""
    splitter = RecursiveCharacterTextSplitter(max_chunk_size=15, max_chunk_overlap=5)
    text = "alpha beta gamma delta epsilon"
    chunks = splitter.split_text(text)
    
    assert len(chunks) >= 2


def test_splitter_preserves_words():
    """测试分割保持单词完整性"""
    splitter = RecursiveCharacterTextSplitter(max_chunk_size=20, max_chunk_overlap=3)
    text = "word1 word2 word3 word4"
    chunks = splitter.split_text(text)
    
    for chunk in chunks:
        words = chunk.split()
        assert all(word in text.split() for word in words)


def test_splitter_single_long_word():
    """测试单个超长单词"""
    splitter = RecursiveCharacterTextSplitter(max_chunk_size=10, max_chunk_overlap=2)
    text = "verylongwordthatexceedschunksize"
    chunks = splitter.split_text(text)
    
    assert len(chunks) >= 1


def test_splitter_empty_text():
    """测试空文本分割"""
    splitter = RecursiveCharacterTextSplitter(max_chunk_size=10, max_chunk_overlap=2)
    chunks = splitter.split_text("")
    
    assert len(chunks) <= 1


def test_splitter_short_text():
    """测试短文本不需分割"""
    splitter = RecursiveCharacterTextSplitter(max_chunk_size=100, max_chunk_overlap=10)
    text = "Short text"
    chunks = splitter.split_text(text)
    
    assert len(chunks) == 1
    assert chunks[0] == text


def test_splitter_multiline_text():
    """测试多行文本分割"""
    splitter = RecursiveCharacterTextSplitter(max_chunk_size=30, max_chunk_overlap=5)
    text = "Line 1 content here\nLine 2 content here\nLine 3 content here"
    chunks = splitter.split_text(text)
    
    assert len(chunks) >= 1


def test_splitter_with_documents(tmp_path):
    """测试对文档列表进行分割"""
    path = tmp_path / "doc.txt"
    content = "a b c d e f g h i j k l m n o p q r s t u v w x y z"
    path.write_text(content, encoding="utf-8")
    
    loader = TextDocLoader(str(path))
    splitter = RecursiveCharacterTextSplitter(max_chunk_size=20, max_chunk_overlap=4)
    
    chunks = loader.load_and_split(splitter)
    
    assert len(chunks) > 1
    for chunk in chunks:
        assert chunk.metadata["source"] == str(path)


def test_splitter_preserves_metadata_across_chunks(tmp_path):
    """测试分割后元数据保持一致"""
    path = tmp_path / "meta_test.txt"
    content = "alpha " * 50
    path.write_text(content, encoding="utf-8")

    loader = TextDocLoader(str(path))
    splitter = RecursiveCharacterTextSplitter(max_chunk_size=30, max_chunk_overlap=5)

    chunks = loader.load_and_split(splitter)

    sources = {chunk.metadata.get("source") for chunk in chunks}
    assert len(sources) == 1
    assert str(path) in sources


# ============================================================================
# MarkdownDocLoader 基础测试
# ============================================================================

def test_markdown_doc_loader_parses_front_matter(tmp_path):
    path = tmp_path / "frontmatter.md"
    path.write_text(
        "---\n"
        "title: Sample Doc\n"
        "tags:\n"
        "  - ai\n"
        "  - markdown\n"
        "---\n\n"
        "# Heading\n"
        "Content here.\n",
        encoding="utf-8",
    )

    loader = MarkdownDocLoader(str(path))
    documents = loader.load()

    assert len(documents) == 1
    doc = documents[0]
    assert doc.page_content.startswith("# Heading")
    assert doc.metadata["source"] == str(path)
    assert doc.metadata["front_matter"]["title"] == "Sample Doc"
    assert doc.metadata["front_matter"]["tags"] == ["ai", "markdown"]


def test_markdown_doc_loader_fetch_from_text():
    loader = MarkdownDocLoader()
    documents = loader.fetch_content(
        {
            "text": "---\ncategory: tests\n---\n\nHello",
            "metadata": {"extra": True},
        }
    )

    doc = documents[0]
    assert doc.page_content.strip() == "Hello"
    assert doc.metadata["front_matter"]["category"] == "tests"
    assert doc.metadata["extra"] is True
    assert doc.metadata["source"] == "inline"


def test_markdown_doc_loader_without_front_matter(tmp_path):
    path = tmp_path / "no_frontmatter.md"
    content = "# Title\n\nPlain text"
    path.write_text(content, encoding="utf-8")

    loader = MarkdownDocLoader(str(path))
    documents = loader.load()

    assert documents[0].page_content == content
    assert "front_matter" not in documents[0].metadata


def test_markdown_doc_loader_requires_source():
    loader = MarkdownDocLoader()
    with pytest.raises(ValueError):
        loader.load()


# ============================================================================
# WordDocLoader 基础测试
# ============================================================================

def _create_docx(path: Path, paragraphs: Iterable[str]) -> None:
    doc = WordDocument()
    for paragraph in paragraphs:
        doc.add_paragraph(paragraph)
    doc.save(path)


def test_word_doc_loader_reads_file(tmp_path):
    docx_path = tmp_path / "sample.docx"
    _create_docx(docx_path, ["Hello", "World"])

    loader = WordDocLoader(str(docx_path))
    documents = loader.load()

    assert len(documents) == 1
    assert "Hello" in documents[0].page_content
    assert "World" in documents[0].page_content
    assert documents[0].metadata["source"] == str(docx_path)


def test_word_doc_loader_fetches_from_bytes(tmp_path):
    docx_path = tmp_path / "bytes.docx"
    _create_docx(docx_path, ["From", "Bytes"])

    payload = docx_path.read_bytes()
    loader = WordDocLoader()
    documents = loader.fetch_content({"bytes": payload, "metadata": {"category": "docx"}})

    assert documents[0].metadata["source"] == "stream"
    assert documents[0].metadata["category"] == "docx"
    assert "From" in documents[0].page_content


def test_word_doc_loader_requires_source():
    loader = WordDocLoader()
    with pytest.raises(ValueError):
        loader.load()


# ============================================================================
# FeishuDocLoader 基础测试
# ============================================================================

def test_feishu_doc_loader_single_document():
    domain = "https://open.feishu.cn"
    responses = {
        (
            "POST",
            f"{domain}/open-apis/auth/v3/tenant_access_token/internal",
        ): {"code": 0, "tenant_access_token": "token", "expire": 7200},
        (
            "GET",
            f"{domain}/open-apis/docx/v1/documents/doc123/raw_content",
        ): {"code": 0, "data": {"content": "Hello Feishu", "title": "Feishu Doc"}},
    }
    session = _StubSession(responses)

    loader = FeishuDocLoader("app", "secret", session=session)
    loader.doc_token = "doc123"

    documents = loader.load()

    assert len(documents) == 1
    doc = documents[0]
    assert doc.page_content == "Hello Feishu"
    assert doc.metadata["doc_token"] == "doc123"
    assert doc.metadata["title"] == "Feishu Doc"
    assert doc.metadata["source"].endswith("doc123")


def test_feishu_doc_loader_space_documents():
    domain = "https://open.feishu.cn"
    responses = {
        (
            "POST",
            f"{domain}/open-apis/auth/v3/tenant_access_token/internal",
        ): {"code": 0, "tenant_access_token": "token", "expire": 7200},
        (
            "GET",
            f"{domain}/open-apis/drive/v1/spaces/space123/nodes",
        ): {
            "code": 0,
            "data": {
                "items": [
                    {
                        "obj_type": "docx",
                        "obj_token": "docA",
                        "node_token": "nodeA",
                        "title": "Doc A",
                    },
                    {
                        "obj_type": "sheet",
                        "obj_token": "sheet1",
                    },
                ],
                "has_more": False,
            },
        },
        (
            "GET",
            f"{domain}/open-apis/docx/v1/documents/docA/raw_content",
        ): {"code": 0, "data": {"content": "Content A"}},
    }
    session = _StubSession(responses)

    loader = FeishuDocLoader("app", "secret", session=session)
    loader.space_id = "space123"
    loader.page_size = 20

    documents = loader.load()

    assert len(documents) == 1
    doc = documents[0]
    assert doc.page_content == "Content A"
    assert doc.metadata["space_id"] == "space123"
    assert doc.metadata["node_token"] == "nodeA"
    assert doc.metadata["doc_type"] == "docx"


def test_feishu_doc_loader_fetch_content_merges_metadata():
    domain = "https://open.feishu.cn"
    responses = {
        (
            "POST",
            f"{domain}/open-apis/auth/v3/tenant_access_token/internal",
        ): {"code": 0, "tenant_access_token": "token", "expire": 7200},
        (
            "GET",
            f"{domain}/open-apis/docx/v1/documents/doc-meta/raw_content",
        ): {"code": 0, "data": {"content": "Meta doc"}},
    }
    session = _StubSession(responses)

    loader = FeishuDocLoader("app", "secret", session=session)
    documents = loader.fetch_content(
        {"doc_token": "doc-meta", "metadata": {"category": "feishu"}}
    )

    assert documents[0].metadata["category"] == "feishu"
    assert documents[0].metadata["doc_token"] == "doc-meta"


def test_feishu_doc_loader_requires_target():
    loader = FeishuDocLoader("app", "secret")
    assert loader.load() == []


# ============================================================================
# PDFDocLoader 基础测试
# ============================================================================

def _create_pdf(tmp_path, name: str, pages: int = 2) -> str:
    pdf_path = tmp_path / name
    writer = PdfWriter()
    for _ in range(pages):
        writer.add_blank_page(width=200, height=200)
    with pdf_path.open("wb") as fh:
        writer.write(fh)
    return str(pdf_path)


def test_pdf_doc_loader_reads_pages(tmp_path):
    pdf_path = _create_pdf(tmp_path, "sample.pdf", pages=3)

    loader = PDFDocLoader(pdf_path)
    documents = loader.load()

    assert len(documents) == 3
    assert {doc.metadata["page_number"] for doc in documents} == {1, 2, 3}
    assert {doc.metadata["source"] for doc in documents} == {pdf_path}


def test_pdf_doc_loader_fetches_from_stream(tmp_path):
    pdf_path = _create_pdf(tmp_path, "stream.pdf")
    with open(pdf_path, "rb") as fh:
        payload = fh.read()

    loader = PDFDocLoader()
    documents = loader.fetch_content({"bytes": payload, "metadata": {"category": "pdf"}})

    assert len(documents) == 2
    for doc in documents:
        assert doc.metadata["source"] == "stream"
        assert doc.metadata["category"] == "pdf"


def test_pdf_doc_loader_requires_source():
    loader = PDFDocLoader()
    with pytest.raises(ValueError):
        loader.load()


def test_pdf_doc_loader_single_page(tmp_path):
    """测试单页PDF"""
    pdf_path = _create_pdf(tmp_path, "single.pdf", pages=1)
    
    loader = PDFDocLoader(pdf_path)
    documents = loader.load()
    
    assert len(documents) == 1
    assert documents[0].metadata["page_number"] == 1


def test_pdf_doc_loader_many_pages(tmp_path):
    """测试多页PDF"""
    pdf_path = _create_pdf(tmp_path, "many.pdf", pages=10)
    
    loader = PDFDocLoader(pdf_path)
    documents = loader.load()
    
    assert len(documents) == 10
    page_numbers = [doc.metadata["page_number"] for doc in documents]
    assert page_numbers == list(range(1, 11))


def test_pdf_doc_loader_metadata_consistency(tmp_path):
    """测试PDF元数据一致性"""
    pdf_path = _create_pdf(tmp_path, "meta.pdf", pages=5)
    
    loader = PDFDocLoader(pdf_path)
    documents = loader.load()
    
    for doc in documents:
        assert "source" in doc.metadata
        assert "page_number" in doc.metadata
        assert doc.metadata["source"] == pdf_path


def test_pdf_doc_loader_nonexistent_file():
    """测试不存在的PDF文件"""
    loader = PDFDocLoader("/nonexistent/file.pdf")
    with pytest.raises((FileNotFoundError, OSError)):
        loader.load()


def test_pdf_doc_loader_with_custom_metadata(tmp_path):
    """测试PDF自定义元数据"""
    pdf_path = _create_pdf(tmp_path, "custom.pdf", pages=2)
    with open(pdf_path, "rb") as fh:
        payload = fh.read()
    
    loader = PDFDocLoader()
    custom_meta = {
        "author": "Test Author",
        "title": "Test Document",
        "category": "technical"
    }
    documents = loader.fetch_content({"bytes": payload, "metadata": custom_meta})
    
    assert len(documents) == 2
    for doc in documents:
        assert doc.metadata["author"] == "Test Author"
        assert doc.metadata["title"] == "Test Document"
        assert doc.metadata["category"] == "technical"


def test_pdf_doc_loader_bytes_empty():
    """测试空字节数据"""
    loader = PDFDocLoader()
    with pytest.raises(Exception):
        loader.fetch_content({"bytes": b""})


# ============================================================================
# 集成测试
# ============================================================================

def test_integration_text_load_split_multiple_files(tmp_path):
    """测试多文件加载和分割的集成"""
    files = []
    for i in range(3):
        path = tmp_path / f"file_{i}.txt"
        content = f"File {i} content: " + "word " * 100
        path.write_text(content, encoding="utf-8")
        files.append(str(path))
    
    splitter = RecursiveCharacterTextSplitter(max_chunk_size=50, max_chunk_overlap=10)
    all_chunks = []
    
    for file_path in files:
        loader = TextDocLoader(file_path)
        chunks = loader.load_and_split(splitter)
        all_chunks.extend(chunks)
    
    assert len(all_chunks) > 3
    sources = {chunk.metadata["source"] for chunk in all_chunks}
    assert len(sources) == 3


def test_integration_mixed_content_types(tmp_path):
    """测试混合内容类型处理"""
    # 文本文件
    txt_path = tmp_path / "text.txt"
    txt_path.write_text("Text content", encoding="utf-8")
    
    # PDF文件
    pdf_path = _create_pdf(tmp_path, "doc.pdf", pages=2)
    
    # 加载文本
    txt_loader = TextDocLoader(str(txt_path))
    txt_docs = txt_loader.load()
    
    # 加载PDF
    pdf_loader = PDFDocLoader(str(pdf_path))
    pdf_docs = pdf_loader.load()
    
    assert len(txt_docs) == 1
    assert len(pdf_docs) == 2
    assert txt_docs[0].metadata["source"] == str(txt_path)
    assert all(doc.metadata["source"] == str(pdf_path) for doc in pdf_docs)


def test_integration_pipeline_with_filtering(tmp_path):
    """测试带过滤的处理管道"""
    path = tmp_path / "pipeline.txt"
    content = "alpha\nbeta\ngamma\ndelta\nepsilon\nzeta\neta\ntheta"
    path.write_text(content, encoding="utf-8")
    
    loader = TextDocLoader(str(path))
    splitter = RecursiveCharacterTextSplitter(max_chunk_size=20, max_chunk_overlap=3)
    chunks = loader.load_and_split(splitter)
    
    # 过滤包含特定内容的块
    filtered = [chunk for chunk in chunks if "alpha" in chunk.page_content or "beta" in chunk.page_content]
    
    assert len(filtered) >= 0
    for chunk in filtered:
        assert "alpha" in chunk.page_content or "beta" in chunk.page_content


def test_integration_document_reconstruction(tmp_path):
    """测试文档重建"""
    path = tmp_path / "reconstruct.txt"
    original_content = "word " * 50
    path.write_text(original_content, encoding="utf-8")
    
    loader = TextDocLoader(str(path))
    splitter = RecursiveCharacterTextSplitter(max_chunk_size=30, max_chunk_overlap=5)
    chunks = loader.load_and_split(splitter)
    
    # 验证所有块都来自同一个源
    assert len(chunks) > 1
    sources = {chunk.metadata["source"] for chunk in chunks}
    assert len(sources) == 1
    assert str(path) in sources


# ============================================================================
# 边界条件和错误处理测试
# ============================================================================

def test_boundary_very_small_chunk_size(tmp_path):
    """测试极小的块大小"""
    path = tmp_path / "small_chunk.txt"
    path.write_text("test content", encoding="utf-8")
    
    loader = TextDocLoader(str(path))
    splitter = RecursiveCharacterTextSplitter(max_chunk_size=1, max_chunk_overlap=0)
    chunks = loader.load_and_split(splitter)
    
    assert len(chunks) >= 1


def test_boundary_zero_overlap(tmp_path):
    """测试零重叠"""
    path = tmp_path / "zero_overlap.txt"
    path.write_text("a b c d e f", encoding="utf-8")
    
    loader = TextDocLoader(str(path))
    splitter = RecursiveCharacterTextSplitter(max_chunk_size=5, max_chunk_overlap=0)
    chunks = loader.load_and_split(splitter)
    
    assert len(chunks) >= 1


def test_boundary_large_overlap(tmp_path):
    """测试大重叠"""
    path = tmp_path / "large_overlap.txt"
    path.write_text("word " * 20, encoding="utf-8")
    
    loader = TextDocLoader(str(path))
    splitter = RecursiveCharacterTextSplitter(max_chunk_size=30, max_chunk_overlap=25)
    chunks = loader.load_and_split(splitter)
    
    assert len(chunks) >= 1


def test_error_invalid_file_path():
    """测试无效文件路径"""
    loader = TextDocLoader(file_path="")
    with pytest.raises((ValueError, FileNotFoundError, OSError)):
        loader.load()


def test_error_directory_instead_of_file(tmp_path):
    """测试传入目录而非文件"""
    loader = TextDocLoader(file_path=str(tmp_path))
    with pytest.raises((IsADirectoryError, PermissionError, OSError)):
        loader.load()


# ============================================================================
# 性能和压力测试
# ============================================================================

def test_performance_large_text_file(tmp_path):
    """测试大文本文件"""
    path = tmp_path / "large.txt"
    content = "line " * 100000
    path.write_text(content, encoding="utf-8")
    
    loader = TextDocLoader(str(path))
    documents = loader.load()
    
    assert len(documents) == 1
    assert len(documents[0].page_content) == 500000


def test_performance_many_small_chunks(tmp_path):
    """测试大量小块分割"""
    path = tmp_path / "many_chunks.txt"
    content = "w " * 10000
    path.write_text(content, encoding="utf-8")
    
    loader = TextDocLoader(str(path))
    splitter = RecursiveCharacterTextSplitter(max_chunk_size=10, max_chunk_overlap=2)
    chunks = loader.load_and_split(splitter)
    
    assert len(chunks) > 100


def test_performance_many_pdf_pages(tmp_path):
    """测试大量PDF页面"""
    pdf_path = _create_pdf(tmp_path, "large.pdf", pages=50)
    
    loader = PDFDocLoader(pdf_path)
    documents = loader.load()
    
    assert len(documents) == 50


# ============================================================================
# 特殊场景测试
# ============================================================================

def test_special_unicode_normalization(tmp_path):
    """测试Unicode规范化"""
    path = tmp_path / "unicode.txt"
    # 同一字符的不同Unicode表示
    content = "café café"  # 一个用组合字符，一个用预组合字符
    path.write_text(content, encoding="utf-8")
    
    loader = TextDocLoader(str(path))
    documents = loader.load()
    
    assert len(documents) == 1


def test_special_line_endings_unix(tmp_path):
    """测试Unix风格换行符"""
    path = tmp_path / "unix.txt"
    content = "line1\nline2\nline3"
    path.write_text(content, encoding="utf-8")
    
    loader = TextDocLoader(str(path))
    documents = loader.load()
    
    assert "\n" in documents[0].page_content


def test_special_line_endings_windows(tmp_path):
    """测试Windows风格换行符"""
    path = tmp_path / "windows.txt"
    content = "line1\r\nline2\r\nline3"
    path.write_bytes(content.encode("utf-8"))
    
    loader = TextDocLoader(str(path))
    documents = loader.load()
    
    assert len(documents) == 1


def test_special_mixed_line_endings(tmp_path):
    """测试混合换行符"""
    path = tmp_path / "mixed.txt"
    content = "line1\nline2\r\nline3\rline4"
    path.write_bytes(content.encode("utf-8"))
    
    loader = TextDocLoader(str(path))
    documents = loader.load()
    
    assert len(documents) == 1


def test_special_tabs_and_spaces(tmp_path):
    """测试制表符和空格混合"""
    path = tmp_path / "tabs.txt"
    content = "word1\tword2  word3\t\tword4    word5"
    path.write_text(content, encoding="utf-8")
    
    loader = TextDocLoader(str(path))
    documents = loader.load()
    
    assert "\t" in documents[0].page_content


def test_special_repeated_separators(tmp_path):
    """测试重复分隔符"""
    path = tmp_path / "repeated.txt"
    content = "word1     word2\n\n\nword3"
    path.write_text(content, encoding="utf-8")
    
    loader = TextDocLoader(str(path))
    documents = loader.load()
    
    assert documents[0].page_content.count(" ") >= 5


# ============================================================================
# 元数据高级测试  
# ============================================================================

def test_metadata_deep_nesting():
    """测试深层嵌套元数据"""
    metadata = {
        "level1": {
            "level2": {
                "level3": {
                    "level4": {"value": "deep"}
                }
            }
        }
    }
    doc = Document(page_content="content", metadata=metadata)
    
    assert doc.metadata["level1"]["level2"]["level3"]["level4"]["value"] == "deep"


def test_metadata_array_operations():
    """测试数组元数据操作"""
    metadata = {"tags": ["tag1", "tag2", "tag3"]}
    doc = Document(page_content="content", metadata=metadata)
    
    doc.metadata["tags"].append("tag4")
    assert len(doc.metadata["tags"]) == 4
    assert "tag4" in doc.metadata["tags"]


def test_metadata_null_values():
    """测试空值元数据"""
    metadata = {"key1": None, "key2": "", "key3": 0, "key4": False}
    doc = Document(page_content="content", metadata=metadata)
    
    assert doc.metadata["key1"] is None
    assert doc.metadata["key2"] == ""
    assert doc.metadata["key3"] == 0
    assert doc.metadata["key4"] is False


# ============================================================================
# ExcelDocLoader Tests
# ============================================================================

def test_excel_doc_loader_reads_file(tmp_path):
    """Test basic Excel file loading with CSV format"""
    excel_path = tmp_path / "test.xlsx"
    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.title = "Sheet1"
    sheet.append(["Name", "Age", "City"])
    sheet.append(["Alice", 30, "New York"])
    sheet.append(["Bob", 25, "London"])
    workbook.save(excel_path)
    
    loader = ExcelDocLoader(str(excel_path), output_format="csv")
    documents = loader.load()
    
    assert len(documents) == 1
    doc = documents[0]
    assert "Name" in doc.page_content
    assert "Alice" in doc.page_content
    assert doc.metadata["sheet_name"] == "Sheet1"
    assert doc.metadata["format"] == "excel"


def test_excel_doc_loader_markdown_format(tmp_path):
    """Test Excel loading with markdown table format"""
    excel_path = tmp_path / "test_md.xlsx"
    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.append(["Product", "Price"])
    sheet.append(["Apple", 1.50])
    sheet.append(["Orange", 2.00])
    workbook.save(excel_path)
    
    loader = ExcelDocLoader(str(excel_path), output_format="markdown")
    documents = loader.load()
    
    assert len(documents) == 1
    assert "| Product | Price |" in documents[0].page_content
    assert "|---|---|" in documents[0].page_content


def test_excel_doc_loader_json_format(tmp_path):
    """Test Excel loading with JSON format"""
    excel_path = tmp_path / "test_json.xlsx"
    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.append(["id", "name"])
    sheet.append([1, "item1"])
    sheet.append([2, "item2"])
    workbook.save(excel_path)
    
    loader = ExcelDocLoader(str(excel_path), output_format="json")
    documents = loader.load()
    
    assert len(documents) == 1
    data = json.loads(documents[0].page_content)
    assert len(data) == 2
    assert data[0]["id"] == "1"


def test_excel_doc_loader_multiple_sheets(tmp_path):
    """Test Excel file with multiple sheets"""
    excel_path = tmp_path / "multi.xlsx"
    workbook = openpyxl.Workbook()
    workbook.remove(workbook.active)
    
    sheet1 = workbook.create_sheet("Sales")
    sheet1.append(["Product", "Sales"])
    sheet1.append(["A", 100])
    
    sheet2 = workbook.create_sheet("Inventory")
    sheet2.append(["Item", "Stock"])
    sheet2.append(["X", 50])
    
    workbook.save(excel_path)
    
    loader = ExcelDocLoader(str(excel_path))
    documents = loader.load()
    
    assert len(documents) == 2
    sheet_names = {doc.metadata["sheet_name"] for doc in documents}
    assert sheet_names == {"Sales", "Inventory"}


def test_excel_doc_loader_raises_when_path_missing():
    """Test Excel loader raises error when no file path provided"""
    loader = ExcelDocLoader()
    with pytest.raises(ValueError):
        loader.load()


def test_excel_doc_loader_fetches_from_metadata_bytes(tmp_path):
    """Test Excel loader with byte stream input"""
    excel_path = tmp_path / "bytes.xlsx"
    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.append(["A", "B"])
    sheet.append([1, 2])
    workbook.save(excel_path)
    
    with open(excel_path, "rb") as f:
        excel_bytes = f.read()
    
    loader = ExcelDocLoader()
    documents = loader.fetch_content({"bytes": excel_bytes})
    
    assert len(documents) == 1
    assert documents[0].metadata["source"] == "stream"


def test_excel_doc_loader_max_rows(tmp_path):
    """Test max_rows parameter limits row processing"""
    excel_path = tmp_path / "many_rows.xlsx"
    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.append(["Col1", "Col2"])
    for i in range(1, 101):
        sheet.append([f"R{i}", f"V{i}"])
    workbook.save(excel_path)
    
    loader = ExcelDocLoader(str(excel_path), max_rows=10)
    documents = loader.load()
    
    assert len(documents) == 1
    line_count = len(documents[0].page_content.split("\n"))
    assert line_count <= 11


# ============================================================================
# PPTDocLoader Tests
# ============================================================================

def test_ppt_doc_loader_reads_file(tmp_path):
    """Test basic PowerPoint file loading"""
    ppt_path = tmp_path / "test.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    
    textbox = slide.shapes.add_textbox(
        Inches(1), Inches(1), Inches(8), Inches(1)
    )
    textbox.text = "Slide Title"
    
    textbox2 = slide.shapes.add_textbox(
        Inches(1), Inches(2), Inches(8), Inches(4)
    )
    textbox2.text = "Slide Content"
    
    presentation.save(str(ppt_path))
    
    loader = PPTDocLoader(str(ppt_path))
    documents = loader.load()
    
    assert len(documents) == 1
    assert documents[0].metadata["slide_number"] == 1
    assert "Slide Title" in documents[0].page_content
    assert "Slide Content" in documents[0].page_content


def test_ppt_doc_loader_multiple_slides(tmp_path):
    """Test PowerPoint with multiple slides"""
    ppt_path = tmp_path / "multi.pptx"
    presentation = Presentation()
    
    for i in range(1, 4):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        textbox = slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(8), Inches(1)
        )
        textbox.text = f"Slide {i}"
    
    presentation.save(str(ppt_path))
    
    loader = PPTDocLoader(str(ppt_path))
    documents = loader.load()
    
    assert len(documents) == 3
    assert documents[0].metadata["slide_number"] == 1
    assert documents[2].metadata["slide_number"] == 3


def test_ppt_doc_loader_with_notes(tmp_path):
    """Test PowerPoint with speaker notes"""
    ppt_path = tmp_path / "notes.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    
    textbox = slide.shapes.add_textbox(
        Inches(1), Inches(1), Inches(8), Inches(1)
    )
    textbox.text = "Main Content"
    
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "Speaker notes here"
    
    presentation.save(str(ppt_path))
    
    loader = PPTDocLoader(str(ppt_path), include_notes=True)
    documents = loader.load()
    
    assert len(documents) == 1
    assert "Main Content" in documents[0].page_content
    assert "Speaker notes here" in documents[0].page_content
    assert documents[0].metadata["has_notes"] is True


def test_ppt_doc_loader_without_notes(tmp_path):
    """Test PowerPoint without including speaker notes"""
    ppt_path = tmp_path / "no_notes.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    
    textbox = slide.shapes.add_textbox(
        Inches(1), Inches(1), Inches(8), Inches(1)
    )
    textbox.text = "Content"
    
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "Hidden notes"
    
    presentation.save(str(ppt_path))
    
    loader = PPTDocLoader(str(ppt_path), include_notes=False)
    documents = loader.load()
    
    assert len(documents) == 1
    assert "Content" in documents[0].page_content
    assert "Hidden notes" not in documents[0].page_content


def test_ppt_doc_loader_raises_when_path_missing():
    """Test PPT loader raises error when no file path provided"""
    loader = PPTDocLoader()
    with pytest.raises(ValueError):
        loader.load()


def test_ppt_doc_loader_fetches_from_metadata_bytes(tmp_path):
    """Test PPT loader with byte stream input"""
    ppt_path = tmp_path / "bytes.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(
        Inches(1), Inches(1), Inches(8), Inches(1)
    )
    textbox.text = "From bytes"
    presentation.save(str(ppt_path))
    
    with open(ppt_path, "rb") as f:
        ppt_bytes = f.read()
    
    loader = PPTDocLoader()
    documents = loader.fetch_content({"bytes": ppt_bytes})
    
    assert len(documents) == 1
    assert documents[0].metadata["source"] == "stream"
    assert "From bytes" in documents[0].page_content


def test_ppt_doc_loader_empty_presentation(tmp_path):
    """Test handling of empty PowerPoint presentation"""
    ppt_path = tmp_path / "empty.pptx"
    presentation = Presentation()
    presentation.save(str(ppt_path))
    
    loader = PPTDocLoader(str(ppt_path))
    documents = loader.load()
    
    assert len(documents) == 0


# ============================================================================
# YuqueDocLoader Tests
# ============================================================================

def test_yuque_doc_loader_single_document():
    """Test loading a single Yuque document"""
    domain = "https://www.yuque.com/api/v2"
    responses = {
        ("GET", f"{domain}/repos/user/repo/docs/test-doc"): {
            "data": {
                "id": 123,
                "slug": "test-doc",
                "title": "Test Document",
                "body": "# Hello Yuque\n\nContent here.",
                "format": "markdown",
                "created_at": "2025-01-01T00:00:00Z",
                "updated_at": "2025-01-02T00:00:00Z",
                "word_count": 100,
                "public": 1,
            }
        },
    }
    session = _StubSession(responses)
    
    loader = YuqueDocLoader(token="test_token", session=session)
    loader.namespace = "user/repo"
    loader.doc_slug = "test-doc"
    
    documents = loader.load()
    
    assert len(documents) == 1
    doc = documents[0]
    assert "Hello Yuque" in doc.page_content
    assert doc.metadata["doc_slug"] == "test-doc"
    assert doc.metadata["title"] == "Test Document"
    assert doc.metadata["namespace"] == "user/repo"


def test_yuque_doc_loader_repository_documents():
    """Test loading all documents from a repository"""
    domain = "https://www.yuque.com/api/v2"
    responses = {
        ("GET", f"{domain}/repos/user/repo/docs"): {
            "data": [
                {"slug": "doc1", "title": "Document 1"},
                {"slug": "doc2", "title": "Document 2"},
            ]
        },
        ("GET", f"{domain}/repos/user/repo/docs/doc1"): {
            "data": {
                "id": 1,
                "slug": "doc1",
                "title": "Document 1",
                "body": "Content of doc1",
                "format": "markdown",
            }
        },
        ("GET", f"{domain}/repos/user/repo/docs/doc2"): {
            "data": {
                "id": 2,
                "slug": "doc2",
                "title": "Document 2",
                "body": "Content of doc2",
                "format": "markdown",
            }
        },
    }
    session = _StubSession(responses)
    
    loader = YuqueDocLoader(token="test_token", session=session)
    loader.namespace = "user/repo"
    
    documents = loader.load()
    
    assert len(documents) == 2
    assert documents[0].metadata["doc_slug"] == "doc1"
    assert documents[1].metadata["doc_slug"] == "doc2"


def test_yuque_doc_loader_fetches_from_metadata():
    """Test Yuque loader with metadata-driven loading"""
    domain = "https://www.yuque.com/api/v2"
    responses = {
        ("GET", f"{domain}/repos/user/repo/docs/meta-doc"): {
            "data": {
                "id": 999,
                "slug": "meta-doc",
                "title": "Meta Document",
                "body": "Metadata test",
                "format": "markdown",
            }
        },
    }
    session = _StubSession(responses)
    
    loader = YuqueDocLoader(token="test_token", session=session)
    documents = loader.fetch_content({
        "namespace": "user/repo",
        "doc_slug": "meta-doc",
        "metadata": {"category": "yuque", "custom": "value"}
    })
    
    assert len(documents) == 1
    assert documents[0].metadata["category"] == "yuque"
    assert documents[0].metadata["custom"] == "value"
    assert documents[0].metadata["doc_slug"] == "meta-doc"


def test_yuque_doc_loader_empty_when_no_target():
    """Test Yuque loader returns empty list when no target specified"""
    loader = YuqueDocLoader(token="test_token")
    documents = loader.load()
    assert documents == []
