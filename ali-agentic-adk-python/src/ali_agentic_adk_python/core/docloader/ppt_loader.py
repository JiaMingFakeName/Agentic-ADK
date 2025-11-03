# Copyright (C) 2025 AIDC-AI
# This project incorporates components from the Open Source Software below.
# The original copyright notices and the licenses under which we received such components are set forth below for informational purposes.
#
# Open Source Software Licensed under the MIT License:
# --------------------------------------------------------------------
# 1. vscode-extension-updater-gitlab 3.0.1 https://www.npmjs.com/package/vscode-extension-updater-gitlab
# Copyright (c) Microsoft Corporation. All rights reserved.
# Copyright (c) 2015 David Owens II
# Copyright (c) Microsoft Corporation.
# Terms of the MIT:
# --------------------------------------------------------------------
# MIT License
#
# Permission is hereby granted, free of charge, to any person obtaining a copy of this software and associated documentation files (the "Software"), to deal in the Software without restriction, including without limitation the rights to use, copy, modify, merge, publish, distribute, sublicense, and/or sell copies of the Software, and to permit persons to whom the Software is furnished to do so, subject to the following conditions:
#
# The above copyright notice and this permission notice shall be included in all copies or substantial portions of the Software.
#
# THE SOFTWARE IS PROVIDED "AS IS", WITHOUT WARRANTY OF ANY KIND, EXPRESS OR IMPLIED, INCLUDING BUT NOT LIMITED TO THE WARRANTIES OF MERCHANTABILITY, FITNESS FOR A PARTICULAR PURPOSE AND NONINFRINGEMENT. IN NO EVENT SHALL THE AUTHORS OR COPYRIGHT HOLDERS BE LIABLE FOR ANY CLAIM, DAMAGES OR OTHER LIABILITY, WHETHER IN AN ACTION OF CONTRACT, TORT OR OTHERWISE, ARISING FROM, OUT OF OR IN CONNECTION WITH THE SOFTWARE OR THE USE OR OTHER DEALINGS IN THE SOFTWARE.

from __future__ import annotations

from io import BytesIO
import logging
from pathlib import Path
from typing import Any, BinaryIO, Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ali_agentic_adk_python.core.docloader.base import BaseLoader
from ali_agentic_adk_python.core.indexes import Document

LOGGER = logging.getLogger(__name__)


class PPTDocLoader(BaseLoader):
    """Load PowerPoint files (PPTX, PPT) and convert slides into ``Document`` objects.
    
    This loader processes PowerPoint presentations and converts each slide into a separate document.
    It extracts:
    - Text from all text boxes and shapes
    - Titles and content placeholders
    - Tables (formatted as markdown or CSV)
    - Notes (speaker notes)
    
    Example:
        >>> loader = PPTDocLoader("presentation.pptx", include_notes=True)
        >>> documents = loader.load()
        >>> for doc in documents:
        ...     print(f"Slide {doc.metadata['slide_number']}")
        ...     print(doc.page_content)
    """

    def __init__(
        self, 
        file_path: Optional[str] = None,
        *,
        include_notes: bool = False,
        table_format: str = "markdown",
        slide_separator: str = "\n\n---\n\n",
    ) -> None:
        """Initialize the PowerPoint document loader.
        
        Args:
            file_path: Path to the PowerPoint file
            include_notes: Whether to include speaker notes
            table_format: Format for tables ('markdown' or 'csv')
            slide_separator: Separator between content sections in a slide
        """
        self.file_path = file_path
        self.include_notes = include_notes
        self.table_format = table_format.lower()
        self.slide_separator = slide_separator
        
        if self.table_format not in ("markdown", "csv"):
            raise ValueError(f"Invalid table_format: {table_format}. Must be 'markdown' or 'csv'")

    def load(self) -> list[Document]:
        """Load the PowerPoint file and convert to documents.
        
        Returns:
            List of Document objects, one per slide
            
        Raises:
            ValueError: If file_path is not set
            FileNotFoundError: If the file does not exist
        """
        if not self.file_path:
            raise ValueError("PPTDocLoader requires `file_path` or metadata-driven loading.")
        return self._load_from_path(self.file_path)

    def fetch_content(self, document_meta: dict[str, Any]) -> list[Document]:
        """Load PowerPoint file using metadata parameters.
        
        Args:
            document_meta: Dictionary containing:
                - file_path or filePath: Path to the file
                - input_stream, stream, or bytes: File content
                - metadata: Additional metadata to merge
                
        Returns:
            List of Document objects
        """
        metadata_hint = document_meta.get("metadata")

        file_path = document_meta.get("file_path") or document_meta.get("filePath")
        if file_path:
            documents = self._load_from_path(str(file_path))
        elif stream := (
            document_meta.get("input_stream")
            or document_meta.get("stream")
            or document_meta.get("binary_stream")
        ):
            documents = self._load_from_stream(stream)
        elif raw_bytes := document_meta.get("bytes"):
            documents = self._load_from_stream(BytesIO(raw_bytes))
        else:
            documents = super().fetch_content(document_meta)

        if metadata_hint:
            for doc in documents:
                doc.metadata.update(metadata_hint)
        return documents

    def _load_from_path(self, file_path: str) -> list[Document]:
        """Load PowerPoint file from a file path.
        
        Args:
            file_path: Path to the PowerPoint file
            
        Returns:
            List of Document objects
            
        Raises:
            FileNotFoundError: If the file does not exist
        """
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        presentation = Presentation(str(path))
        return self._extract_documents(presentation, str(path))

    def _load_from_stream(self, stream: BinaryIO | BytesIO) -> list[Document]:
        """Load PowerPoint file from a byte stream.
        
        Args:
            stream: Binary stream containing PowerPoint file data
            
        Returns:
            List of Document objects
        """
        presentation = Presentation(stream)
        return self._extract_documents(presentation, None)

    def _extract_documents(
        self, 
        presentation: Presentation, 
        source: Optional[str]
    ) -> list[Document]:
        """Extract documents from a PowerPoint presentation.
        
        Args:
            presentation: Opened PowerPoint presentation
            source: Source file path (None if from stream)
            
        Returns:
            List of Document objects
        """
        documents: list[Document] = []
        
        for slide_idx, slide in enumerate(presentation.slides, start=1):
            try:
                content_parts = []
                
                # Extract text from all shapes
                for shape in slide.shapes:
                    shape_text = self._extract_shape_text(shape)
                    if shape_text:
                        content_parts.append(shape_text)
                
                # Extract speaker notes if requested
                notes_text = ""
                if self.include_notes and slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    notes_text = notes_slide.notes_text_frame.text if notes_slide.notes_text_frame else ""
                    if notes_text.strip():
                        content_parts.append(f"**Notes:**\n{notes_text}")
                
                # Combine all content
                content = self.slide_separator.join(content_parts)
                
                metadata = {
                    "source": source or "stream",
                    "slide_number": slide_idx,
                    "format": "powerpoint",
                    "has_notes": bool(notes_text),
                    "shape_count": len(slide.shapes),
                }
                
                document = Document(page_content=content, metadata=metadata)
                documents.append(document)
                
            except Exception as exc:
                LOGGER.warning(f"Failed to process slide {slide_idx}: {exc}")
                continue
        
        return documents

    def _extract_shape_text(self, shape) -> str:
        """Extract text from a PowerPoint shape.
        
        Args:
            shape: PowerPoint shape object
            
        Returns:
            Extracted text content
        """
        text_parts = []
        
        # Check if shape has text frame
        if hasattr(shape, "text") and shape.text:
            text_parts.append(shape.text)
        
        # Handle tables specially
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table_text = self._extract_table_text(shape.table)
            if table_text:
                text_parts.append(table_text)
        
        # Handle grouped shapes
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for grouped_shape in shape.shapes:
                grouped_text = self._extract_shape_text(grouped_shape)
                if grouped_text:
                    text_parts.append(grouped_text)
        
        # Handle text boxes and placeholders
        elif hasattr(shape, "text_frame") and shape.text_frame:
            frame_text = []
            for paragraph in shape.text_frame.paragraphs:
                para_text = paragraph.text.strip()
                if para_text:
                    frame_text.append(para_text)
            if frame_text:
                text_parts.append("\n".join(frame_text))
        
        return "\n".join(text_parts)

    def _extract_table_text(self, table) -> str:
        """Extract and format text from a PowerPoint table.
        
        Args:
            table: PowerPoint table object
            
        Returns:
            Formatted table text
        """
        rows = []
        
        for row in table.rows:
            row_cells = []
            for cell in row.cells:
                cell_text = cell.text.strip()
                row_cells.append(cell_text)
            rows.append(row_cells)
        
        if not rows:
            return ""
        
        if self.table_format == "markdown":
            return self._format_table_as_markdown(rows)
        else:
            return self._format_table_as_csv(rows)

    def _format_table_as_markdown(self, rows: list[list[str]]) -> str:
        """Format table rows as Markdown.
        
        Args:
            rows: List of rows, each row is a list of cell values
            
        Returns:
            Markdown table string
        """
        if not rows:
            return ""
        
        lines = []
        
        # First row as header
        if rows:
            header = rows[0]
            lines.append("| " + " | ".join(header) + " |")
            lines.append("|" + "|".join(["---" for _ in header]) + "|")
        
        # Data rows
        for row in rows[1:]:
            # Escape pipe characters
            escaped_row = [cell.replace("|", chr(92) + "|") for cell in row]
            lines.append("| " + " | ".join(escaped_row) + " |")
        
        return "\n".join(lines)

    def _format_table_as_csv(self, rows: list[list[str]]) -> str:
        """Format table rows as CSV.
        
        Args:
            rows: List of rows, each row is a list of cell values
            
        Returns:
            CSV string
        """
        lines = []
        for row in rows:
            # Escape commas and quotes
            escaped_row = []
            for cell in row:
                if "," in cell or '"' in cell:
                    escaped_cell = f'"{cell.replace(chr(34), chr(34) + chr(34))}"'
                    escaped_row.append(escaped_cell)
                else:
                    escaped_row.append(cell)
            lines.append(",".join(escaped_row))
        return "\n".join(lines)


